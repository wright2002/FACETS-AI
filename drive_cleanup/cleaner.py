import json
import os
import magic
from collections import defaultdict
from docx import Document
import fitz  # PyMuPDF
import openpyxl
from pptx import Presentation
from sklearn.cluster import KMeans
from sentence_transformers import SentenceTransformer
from sklearn.decomposition import PCA
from sklearn.feature_extraction.text import TfidfVectorizer


def find_duplicates(file_metadata_list):
    hash_map = defaultdict(list)
    for file_info in file_metadata_list:
        file_hash = file_info.get("hash_sha256")
        if file_hash and not file_hash.startswith("ERROR"):
            hash_map[file_hash].append(file_info["path"])
    duplicates = {h: paths for h, paths in hash_map.items() if len(paths) > 1}
    return duplicates

def extension_matches_mime(ext, mime):
    expected_mimes = {
        "pdf": "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "png": "image/png",
        "gif": "image/gif"
    }
    expected = expected_mimes.get(ext)
    return expected == mime

def detect_obfuscated_types(file_metadata_list):
    suspicious_files = []
    for file_info in file_metadata_list:
        file_path = file_info.get("path")
        _, ext = os.path.splitext(file_path)
        ext = ext.lower().strip(".")

        try:
            mime = magic.Magic(mime=True).from_file(file_path)
            file_info["detected_mime"] = mime
            if not extension_matches_mime(ext, mime):
                file_info["suspicious"] = True
                suspicious_files.append({
                    "path": file_path,
                    "extension": ext,
                    "detected_mime": mime
                })
            else:
                file_info["suspicious"] = False
        except Exception as e:
            file_info["mime_error"] = str(e)

    return suspicious_files

def extract_text_from_file(file_path):
    _, ext = os.path.splitext(file_path)
    ext = ext.lower().strip(".")

    try:
        if ext == "docx":
            doc = Document(file_path)
            return "\n".join([p.text for p in doc.paragraphs])
        elif ext == "pdf":
            doc = fitz.open(file_path)
            return "\n".join([page.get_text() for page in doc])
        elif ext == "xlsx":
            wb = openpyxl.load_workbook(file_path, data_only=True)
            text = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    text.append(" ".join([str(cell) if cell is not None else "" for cell in row]))
            return "\n".join(text)
        elif ext == "pptx":
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        else:
            return ""
    except Exception as e:
        return f"ERROR: {e}"

def extract_text_dataset(file_metadata):
    dataset = []
    for entry in file_metadata:
        file_path = entry.get("path")
        extracted_text = extract_text_from_file(file_path)
        if extracted_text and not extracted_text.startswith("ERROR"):
            dataset.append({
                "path": file_path,
                "text": extracted_text
            })
    return dataset

def cluster_documents(text_dataset, n_clusters=4):
    texts = [d["text"] for d in text_dataset]
    paths = [d["path"] for d in text_dataset]

    # Embed using sentence-transformers
    model = SentenceTransformer("all-MiniLM-L6-v2")
    embeddings = model.encode(texts, show_progress_bar=True)

    # Reduce dimensions with PCA
    n_components = min(25, len(embeddings), len(embeddings[0]))
    embeddings = PCA(n_components=n_components).fit_transform(embeddings)

    # KMeans clustering
    kmeans = KMeans(n_clusters=n_clusters, random_state=42)
    labels = kmeans.fit_predict(embeddings)

    # TF-IDF for naming clusters by keywords
    tfidf = TfidfVectorizer(stop_words='english', max_features=1000)
    tfidf_matrix = tfidf.fit_transform(texts)
    terms = tfidf.get_feature_names_out()

    cluster_names = {}
    for i in range(n_clusters):
        indices = [j for j, label in enumerate(labels) if label == i]
        if not indices:
            cluster_names[i] = f"cluster_{i}"
            continue

        cluster_text = tfidf_matrix[indices].mean(axis=0).A1
        top_term_indices = cluster_text.argsort()[-3:][::-1]
        top_keywords = [terms[idx] for idx in top_term_indices]
        cluster_names[i] = "_".join(top_keywords) if top_keywords else f"cluster_{i}"

    # Build recommendations
    recommendations = []
    for path, label in zip(paths, labels):
        folder_name = cluster_names[label]
        recommended_folder = f"D:\\organized\\{folder_name}\\"
        recommendations.append({
            "original_path": path,
            "cluster_id": int(label),
            "suggested_folder": recommended_folder,
            "new_path": os.path.join(recommended_folder, os.path.basename(path))
        })

    return recommendations

def generate_cleanup_report(file_metadata_path, report_output_path, rec_output_path):
    with open(file_metadata_path, "r") as f:
        file_metadata = json.load(f)

    duplicates = find_duplicates(file_metadata)
    obfuscated = detect_obfuscated_types(file_metadata)
    text_dataset = extract_text_dataset(file_metadata)
    move_recommendations = cluster_documents(text_dataset)

    path_map = {rec["original_path"]: rec for rec in move_recommendations}

    report = {
        "duplicate_files": [],
        "obfuscated_files": obfuscated,
        "restructured_files": move_recommendations,
        "summary": {
            "total_files_scanned": len(file_metadata),
            "total_duplicate_sets": len(duplicates),
            "total_duplicate_files": sum(len(paths) for paths in duplicates.values()),
            "total_suspicious_files": len(obfuscated)
        }
    }

    for file_hash, paths in duplicates.items():
        report["duplicate_files"].append({
            "hash_sha256": file_hash,
            "files": paths,
            "keep": paths[0],
            "remove": paths[1:]
        })

    with open(report_output_path, "w", encoding="utf-8") as f:
        json.dump(report, f, indent=2)

    with open(rec_output_path, "w", encoding="utf-8") as f:
        json.dump(move_recommendations, f, indent=2)

    return report_output_path, rec_output_path

def main():
    generate_cleanup_report(
        "scan_results_files.json",
        "cleanup_report.json",
        "move_recommendations.json"
    )

main()
